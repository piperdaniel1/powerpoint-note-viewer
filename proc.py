import collections
import collections.abc
from pptx import Presentation
import sys
import os

def get_pptx_file():
    dir_items = os.listdir()

    for item in dir_items:
        if os.path.isfile(item):
            if item.split(".")[-1] == "pptx":
                return item

    return None

def main():
    INPUT_FILE = get_pptx_file()
    OUTPUT_FILE = "parsed_powerpoint_notes.dispnt"

    if len(sys.argv) > 5 or len(sys.argv) == 2 or len(sys.argv) == 4:
        print("Error, invalid number of args provided.") 
        print("Format: python3 proc.py [--input <input_file> (optional)] [--output <output_file> (optional, .dispnt format)]")
        return

    if len(sys.argv) >= 3:
        if sys.argv[1] == "--input":
            INPUT_FILE = sys.argv[2]
        elif sys.argv[1] == "--output":
            OUTPUT_FILE = sys.argv[2]
        else:
            print("Error, invalid arg flag provided.") 
            print("Format: python3 proc.py [--input <input_file> (optional)] [--output <output_file> (optional)]")

    if len(sys.argv) == 5:
        if sys.argv[3] == "--input":
            INPUT_FILE = sys.argv[4]
        elif sys.argv[3] == "--output":
            OUTPUT_FILE = sys.argv[4]
        else:
            print("Error, invalid arg flag provided.") 
            print("Format: python3 proc.py [--input <input_file> (optional)] [--output <output_file> (optional)]")

    if INPUT_FILE == None:
        print("Error, no .pptx input file found in current directory. Specify one manually or put a .pptx file in the current working directory.")
        print("Format: python3 proc.py [--input <input_file> (optional)] [--output <output_file> (optional)]")
        return

    with open(OUTPUT_FILE, "w") as output:
        pres = Presentation(INPUT_FILE)
        for slide in pres.slides:
            output.write("\n<!--DIV-->\n")
            if slide.has_notes_slide:
                output.write(slide.notes_slide.notes_text_frame.text)
            else:
                output.write(" -- No notes provided -- ")

main()
